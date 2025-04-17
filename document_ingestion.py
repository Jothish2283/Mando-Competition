import io
import json
import shutil
import os

import pdfplumber
import docx
import pptx
import pandas as pd
from PIL import Image

# Try importing OCR backends:
_TESSERACT_AVAILABLE = shutil.which("tesseract") is not None
try:
    from paddleocr import PaddleOCR
    import numpy as np
    PADDLE_AVAILABLE = True
    _paddle_ocr = PaddleOCR(lang="en", use_gpu=False)
except ImportError:
    PADDLE_AVAILABLE = False

try:
    from google.cloud import vision
    GCV_AVAILABLE = True
    _gcv_client = vision.ImageAnnotatorClient()
except ImportError:
    GCV_AVAILABLE = False

def _ocr_image_bytes(image_bytes: bytes) -> str:
    """
    Given raw image bytes, run OCR via:
     1) PaddleOCR (if installed)
     2) Google Vision (if installed & creds set)
     3) Tesseract (if on PATH)
    """
    # PaddleOCR
    if PADDLE_AVAILABLE:
        # PaddleOCR expects a numpy array
        arr = np.array(Image.open(io.BytesIO(image_bytes)))
        result = _paddle_ocr.ocr(arr, cls=False)
        # result is list of [ [ (x1...), (text, conf) ], ... ]
        lines = []
        for block in result:
            for line in block:
                lines.append(line[1][0])
        return "\n".join(lines)

    # Google Vision
    if GCV_AVAILABLE and os.getenv("GOOGLE_APPLICATION_CREDENTIALS"):
        try:
            img = vision.Image(content=image_bytes)
            resp = _gcv_client.document_text_detection(image=img)
            return resp.full_text_annotation.text or ""
        except Exception as e:
            return f"[GCV Error] {e}"

    # Tesseract
    if _TESSERACT_AVAILABLE:
        try:
            import pytesseract
            img = Image.open(io.BytesIO(image_bytes))
            return pytesseract.image_to_string(img)
        except Exception as e:
            return f"[Tesseract OCR Error] {e}"

    return "[OCR Skipped: no OCR engine available]"


def ingest_pdf(file_bytes: bytes) -> str:
    """
    Extract text from PDF; for pages with no text, fallback to _ocr_image_bytes.
    """
    texts = []
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                txt = page.extract_text() or ""
                if txt.strip():
                    texts.append(txt)
                else:
                    # render page to image
                    pil_img = page.to_image(resolution=150).original
                    buf = io.BytesIO()
                    pil_img.save(buf, format="PNG")
                    texts.append(_ocr_image_bytes(buf.getvalue()))
    except Exception as e:
        texts.append(f"[PDF Open Error] {e}")
    return "\n".join(texts)


def ingest_docx(file_bytes: bytes) -> str:
    """Extract text from a .docx."""
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[DOCX Error] {e}"


def ingest_pptx(file_bytes: bytes) -> str:
    """Extract text from a .pptx."""
    try:
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, "text") and shp.text.strip():
                    texts.append(shp.text)
        return "\n".join(texts)
    except Exception as e:
        return f"[PPTX Error] {e}"


def ingest_csv(file_bytes: bytes):
    """Load CSV into DataFrame and return (df, csv_text)."""
    try:
        df = pd.read_csv(io.BytesIO(file_bytes))
        return df, df.to_csv(index=False)
    except Exception as e:
        return None, f"[CSV Error] {e}"


def ingest_excel(file_bytes: bytes):
    """Load Excel into DataFrame and return (df, csv_text)."""
    try:
        df = pd.read_excel(io.BytesIO(file_bytes))
        return df, df.to_csv(index=False)
    except Exception as e:
        return None, f"[Excel Error] {e}"


def ingest_json(file_bytes: bytes) -> str:
    """Pretty‑print JSON."""
    try:
        data = json.load(io.BytesIO(file_bytes))
        return json.dumps(data, indent=2)
    except Exception as e:
        return f"[JSON Error] {e}"


def ingest_txt(file_bytes: bytes) -> str:
    """Read a plaintext file."""
    try:
        return file_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        return f"[TXT Error] {e}"


def ingest_image(file_bytes: bytes) -> str:
    """Run OCR on an image file."""
    return _ocr_image_bytes(file_bytes)
